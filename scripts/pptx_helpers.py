"""
pptx_helpers.py — Reusable PowerPoint building blocks for Noah Askin's decks.

TWO STYLE FAMILIES:

  NSASlidesV2 (CURRENT — use for new decks)
    Gray background (#4D4D4D), green 26pt sentence-case titles, large
    centered statements alternating white/green, green rounded callout
    boxes, light striped tables, generous whitespace, minimal cards.
    Matches: 260319_NSA_Berkeley Culture.pptx, 260702 Nagymaros deck.

    from pptx_helpers import NSASlidesV2
    deck = NSASlidesV2("My Deck Title")
    s = deck.new_slide()
    deck.title(s, "Slide title in sentence case")
    deck.statements(s, [
        {"text": "A big centered statement", "size": 22},
        {"text": "The green follow-up", "color": deck.C.GREEN},
    ], top=Inches(2))
    deck.save("output.pptx")

  NSASlides (LEGACY — only for editing older decks)
    Dark background (#1E1E1E), white bold titles, dark cards (#333333).

Both: 10" x 6.25" slides, Arial throughout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn


class Colors:
    """LEGACY palette (dark family — older decks only)."""
    BG         = RGBColor(0x1E, 0x1E, 0x1E)  # Dark background
    GREEN      = RGBColor(0x92, 0xD0, 0x50)  # Primary accent (lime green)
    GREEN2     = RGBColor(0x85, 0xBD, 0x41)  # Secondary green (muted)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)  # Body text
    LIGHT      = RGBColor(0xCC, 0xCC, 0xCC)  # Muted text
    DIM        = RGBColor(0x88, 0x88, 0x88)  # Dimmer text / footnotes
    DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)  # Card / table cell backgrounds
    MID_GRAY   = RGBColor(0x4D, 0x4D, 0x4D)  # Borders, separators
    PLACEHOLDER = RGBColor(0x2A, 0x2A, 0x2A) # Figure placeholder background
    RED        = RGBColor(0xFF, 0x44, 0x44)  # Negative / attention


FONT = "Arial"
SLIDE_W = Inches(10)
SLIDE_H = Inches(6.25)


class NSASlides:
    """LEGACY builder (dark family). Use NSASlidesV2 for new decks."""

    C = Colors  # Expose palette as class attribute for external use

    def __init__(self, title="Untitled", author="Noah Askin"):
        self.pres = Presentation()
        self.pres.slide_width = SLIDE_W
        self.pres.slide_height = SLIDE_H
        self.pres.core_properties.author = author
        self.pres.core_properties.title = title
        self._blank = self.pres.slide_layouts[6]

    # ── Background & shapes ──────────────────────────────────

    def add_bg(self, slide, color=None):
        """Set solid background color on a slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color or Colors.BG

    def add_rect(self, slide, left, top, width, height, fill_color):
        """Add a borderless filled rectangle."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    # ── Text helpers ─────────────────────────────────────────

    def add_text(self, slide, left, top, width, height, text,
                 size=20, color=None, bold=False, italic=False,
                 align=PP_ALIGN.LEFT, font=FONT):
        """Add a simple text box."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.name = font
        p.font.color.rgb = color or Colors.WHITE
        p.font.bold = bold
        p.font.italic = italic
        p.alignment = align
        return txBox

    def add_bullets(self, slide, left, top, width, height, items,
                    size=18, color=None, spacing=Pt(6)):
        """
        Add bulleted list. items: list of str or dict with keys:
        text, color, bold, italic, size.
        """
        color = color or Colors.WHITE
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            if isinstance(item, str):
                p.text = item
                p.font.color.rgb = color
                p.font.size = Pt(size)
            else:
                run = p.add_run()
                run.text = item.get("text", "")
                run.font.color.rgb = item.get("color", color)
                run.font.bold = item.get("bold", False)
                run.font.italic = item.get("italic", False)
                run.font.size = Pt(item.get("size", size))

            p.font.name = FONT
            p.space_after = spacing
            pPr = p._p.get_or_add_pPr()
            buChar = etree.SubElement(pPr, qn('a:buChar'))
            buChar.set('char', '\u2022')

        return txBox

    def add_green_header_bullets(self, slide, left, top, width, height, sections):
        """
        Sectioned list with green headers and white bullet sub-items.
        sections: list of (header_text, [bullet_items])
        """
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        first = True
        for header, bullets in sections:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = header
            p.font.size = Pt(20)
            p.font.name = FONT
            p.font.color.rgb = Colors.GREEN
            p.font.bold = True
            p.space_before = Pt(12)
            p.space_after = Pt(4)

            for b in bullets:
                p = tf.add_paragraph()
                if isinstance(b, str):
                    p.text = b
                    p.font.color.rgb = Colors.WHITE
                    p.font.size = Pt(16)
                else:
                    run = p.add_run()
                    run.text = b.get("text", "")
                    run.font.color.rgb = b.get("color", Colors.WHITE)
                    run.font.size = Pt(b.get("size", 16))
                    run.font.bold = b.get("bold", False)
                p.font.name = FONT
                p.space_after = Pt(3)
                pPr = p._p.get_or_add_pPr()
                buChar = etree.SubElement(pPr, qn('a:buChar'))
                buChar.set('char', '\u2022')
                pPr.set('marL', str(int(Pt(18))))
                pPr.set('indent', str(int(Pt(-14))))

        return txBox

    # ── Figure placeholders ──────────────────────────────────

    def figure_placeholder(self, slide, left, top, width, height, label):
        """Gray box with centered label — drop in the real figure later."""
        self.add_rect(slide, left, top, width, height, Colors.PLACEHOLDER)
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        border.fill.background()
        border.line.color.rgb = Colors.MID_GRAY
        border.line.width = Pt(1)
        self.add_text(
            slide, left, top + height // 2 - Inches(0.25), width, Inches(0.5),
            f"[ INSERT: {label} ]",
            size=14, color=Colors.DIM, align=PP_ALIGN.CENTER
        )

    # ── Slide templates ──────────────────────────────────────

    def new_slide(self, title="", subtitle_text=None):
        """Standard content slide with white title at top."""
        slide = self.pres.slides.add_slide(self._blank)
        self.add_bg(slide)
        if title:
            self.add_text(slide, Inches(0.5), Inches(0.15), Inches(9), Inches(0.8),
                          title, size=28, color=Colors.WHITE, bold=True)
        if subtitle_text:
            self.add_text(slide, Inches(0.5), Inches(0.75), Inches(9), Inches(0.4),
                          subtitle_text, size=16, color=Colors.DIM)
        return slide

    def section_slide(self, title, subtitle=""):
        """Section divider with green accent bar and large green title."""
        slide = self.pres.slides.add_slide(self._blank)
        self.add_bg(slide)
        self.add_rect(slide, Inches(0.5), Inches(2.6), Inches(1.2), Inches(0.05),
                      Colors.GREEN)
        self.add_text(slide, Inches(0.5), Inches(2.8), Inches(9), Inches(1.2),
                      title, size=36, color=Colors.GREEN, bold=True)
        if subtitle:
            self.add_text(slide, Inches(0.5), Inches(3.9), Inches(9), Inches(0.6),
                          subtitle, size=18, color=Colors.LIGHT)
        return slide

    def title_slide(self, title, authors, accent_bar_width=Inches(2)):
        """
        Opening title slide.
        authors: list of str (name + affiliation lines)
        """
        slide = self.pres.slides.add_slide(self._blank)
        self.add_bg(slide)
        self.add_text(slide, Inches(1.0), Inches(1.2), Inches(7), Inches(1.5),
                      title, size=30, color=Colors.GREEN, bold=True)
        self.add_rect(slide, Inches(1.0), Inches(2.7), accent_bar_width,
                      Inches(0.04), Colors.GREEN)
        for i, author in enumerate(authors):
            self.add_text(slide, Inches(1.0), Inches(3.2 + i * 0.4),
                          Inches(7), Inches(0.4),
                          author, size=16, color=Colors.WHITE, bold=True)
        return slide

    def thank_you_slide(self, email, coauthors=None):
        """Closing slide with email and optional coauthor list."""
        slide = self.pres.slides.add_slide(self._blank)
        self.add_bg(slide)
        self.add_text(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1.0),
                      "Thank you!", size=40, color=Colors.GREEN, bold=True,
                      align=PP_ALIGN.CENTER)
        self.add_rect(slide, Inches(4.0), Inches(2.5), Inches(2.0),
                      Inches(0.04), Colors.GREEN)
        self.add_text(slide, Inches(0.5), Inches(3.0), Inches(9), Inches(0.5),
                      email, size=20, color=Colors.WHITE, align=PP_ALIGN.CENTER)
        if coauthors:
            for i, ca in enumerate(coauthors):
                self.add_text(slide, Inches(0.5), Inches(4.0 + i * 0.4),
                              Inches(9), Inches(0.4),
                              ca, size=16, color=Colors.LIGHT,
                              align=PP_ALIGN.CENTER)
        return slide

    # ── Table helpers ────────────────────────────────────────

    def add_striped_table(self, slide, left, top, col_widths, headers, rows,
                          header_size=13, row_size=12):
        """
        Two-column (or more) striped table with green header bar.
        col_widths: list of Inches values
        headers: list of str
        rows: list of tuples matching len(headers)
        """
        x = left
        for j, (hdr, w) in enumerate(zip(headers, col_widths)):
            self.add_rect(slide, x, top, w, Inches(0.35), Colors.GREEN)
            self.add_text(slide, x + Inches(0.1), top, w - Inches(0.2),
                          Inches(0.35), hdr, size=header_size,
                          color=Colors.BG, bold=True)
            x += w + Inches(0.05)

        for i, row in enumerate(rows):
            y = top + Inches(0.40 + i * 0.43)
            bg = Colors.DARK_GRAY if i % 2 == 0 else Colors.BG
            x = left
            for j, (cell, w) in enumerate(zip(row, col_widths)):
                self.add_rect(slide, x, y, w, Inches(0.4), bg)
                clr = Colors.GREEN if j == 0 else Colors.LIGHT
                self.add_text(slide, x + Inches(0.1), y, w - Inches(0.2),
                              Inches(0.4), cell, size=row_size,
                              color=clr, bold=(j == 0))
                x += w + Inches(0.05)

    # ── Save ─────────────────────────────────────────────────

    def save(self, path):
        """Save the presentation and print summary."""
        self.pres.save(path)
        n = len(self.pres.slides)
        print(f"Saved {n} slides to: {path}")
        return path


# ═══════════════════════════════════════════════════════════════
# CURRENT style family — gray bg, green titles, centered statements
# (matches 260319_NSA_Berkeley Culture.pptx / 260702 Nagymaros deck)
# ═══════════════════════════════════════════════════════════════

class ColorsV2:
    """Palette for the current deck family."""
    BG        = RGBColor(0x4D, 0x4D, 0x4D)  # warm medium gray background
    GREEN     = RGBColor(0x92, 0xD0, 0x50)  # titles, subheads, emphasis
    GREEN_BOX = RGBColor(0x85, 0xBD, 0x41)  # callout box fill
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)  # primary body text
    LIGHT     = RGBColor(0xE0, 0xE0, 0xE0)  # muted white
    DIM       = RGBColor(0xBD, 0xBD, 0xBD)  # footnotes / credits / captions
    TBL_HDR_T = RGBColor(0x40, 0x40, 0x40)  # dark text on green table header
    TBL_TXT   = RGBColor(0x33, 0x33, 0x33)  # dark text in light table rows
    ROW_A     = RGBColor(0xE7, 0xEF, 0xDA)  # pale green table stripe
    ROW_B     = RGBColor(0xF6, 0xF8, 0xF1)  # near-white table stripe
    PH_BG     = RGBColor(0xEF, 0xEF, 0xEF)  # figure placeholder background
    PH_TXT    = RGBColor(0x8A, 0x8A, 0x8A)  # figure placeholder label


class NSASlidesV2:
    """Builder for the CURRENT presentation style.

    Design language:
      - Background #4D4D4D; Arial; 10" x 6.25"
      - Green 26pt sentence-case titles, top-left
      - Large centered statements alternating white/green (statements())
      - Green bold subheads + bullets/numbered steps for columns
      - Green rounded callout boxes for punchlines (callout())
      - Light striped tables: green header, pale rows, dark text
      - Light figure placeholders (real figures are white R plots)

    Rich text: content items are dicts, either {"text", "color", "size",
    "bold", "italic", "gap"} or {"runs": [(text, color, size, bold,
    italic), ...], "gap": Pt(n)} for mixed formatting within a line.
    """

    C = ColorsV2

    def __init__(self, title="Untitled", author="Noah Askin"):
        self.pres = Presentation()
        self.pres.slide_width = SLIDE_W
        self.pres.slide_height = SLIDE_H
        self.pres.core_properties.author = author
        self.pres.core_properties.title = title
        self._blank = self.pres.slide_layouts[6]

    # ── Core primitives ──────────────────────────────────────

    def new_slide(self):
        """Blank slide with the gray background. Add a title with title()."""
        s = self.pres.slides.add_slide(self._blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = ColorsV2.BG
        return s

    @staticmethod
    def notes(slide, text):
        """Set speaker notes."""
        slide.notes_slide.notes_text_frame.text = text

    @staticmethod
    def run(r, text, color=None, size=20, bold=False, italic=False):
        """Format a run (pass p.add_run())."""
        r.text = text
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color or ColorsV2.WHITE
        r.font.bold = bold
        r.font.italic = italic
        return r

    @staticmethod
    def para(tf, first):
        """First paragraph of a text frame, or a new one."""
        return tf.paragraphs[0] if first else tf.add_paragraph()

    @staticmethod
    def textbox(slide, x, y, w, h):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.text_frame.word_wrap = True
        return tb

    def _items_to_runs(self, it):
        return it.get("runs") or [(it.get("text", ""),
                                   it.get("color", ColorsV2.WHITE),
                                   it.get("size", None),
                                   it.get("bold", False),
                                   it.get("italic", False))]

    # ── Layout blocks ────────────────────────────────────────

    def title(self, slide, text):
        """Green 26pt sentence-case title, top-left."""
        tb = self.textbox(slide, Inches(0.55), Inches(0.26),
                          Inches(8.9), Inches(0.75))
        p = tb.text_frame.paragraphs[0]
        self.run(p.add_run(), text, ColorsV2.GREEN, 26)
        return tb

    def statements(self, slide, items, top, left=Inches(0.8),
                   width=Inches(8.4), default_size=20):
        """Centered stacked statements — the family's core pattern.
        Alternate white/green across items for rhythm."""
        tb = self.textbox(slide, left, top, width, Inches(0.5))
        tf = tb.text_frame
        tf.auto_size = None
        for i, it in enumerate(items):
            p = self.para(tf, i == 0)
            p.alignment = PP_ALIGN.CENTER
            for t, c, sz, b, ital in self._items_to_runs(it):
                self.run(p.add_run(), t, c, sz or default_size, b, ital)
            p.space_after = it.get("gap", Pt(18))
        return tb

    def bullets(self, slide, x, y, w, h, items, size=17, gap=Pt(8),
                numbered=False):
        """Left-aligned bulleted or numbered list."""
        tb = self.textbox(slide, x, y, w, h)
        tf = tb.text_frame
        for i, it in enumerate(items):
            p = self.para(tf, i == 0)
            for t, c, sz, b, ital in self._items_to_runs(it):
                self.run(p.add_run(), t, c, sz or size, b, ital)
            p.space_after = it.get("gap", gap)
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(int(Pt(20))))
            pPr.set('indent', str(int(Pt(-20))))
            etree.SubElement(pPr, qn('a:buFont'), {'typeface': FONT})
            if numbered:
                etree.SubElement(pPr, qn('a:buAutoNum'),
                                 {'type': 'arabicPeriod'})
            else:
                etree.SubElement(pPr, qn('a:buChar'), {'char': '•'})
        return tb

    def subhead(self, slide, x, y, w, text, size=20):
        """Green bold subhead for column layouts."""
        tb = self.textbox(slide, x, y, w, Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        self.run(p.add_run(), text, ColorsV2.GREEN, size, bold=True)
        return tb

    def callout(self, slide, text, y, x=Inches(1.1), w=Inches(7.8),
                h=Inches(0.95), size=18):
        """Green rounded box, white bold-italic centered — the punchline.
        Widen (x=0.7, w=8.6) if the text wraps awkwardly."""
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = ColorsV2.GREEN_BOX
        box.line.fill.background()
        box.adjustments[0] = 0.12
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.25)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self.run(p.add_run(), text, ColorsV2.WHITE, size, bold=True,
                 italic=True)
        return box

    def light_table(self, slide, x, y, col_w, headers, rows, hdr_size=14,
                    row_size=13, row_h=Inches(0.42), first_col_bold=True):
        """Light striped table: green header, pale rows, dark text.
        Returns the y-coordinate just below the table.
        NOTE: leave y >= 1.5" when the slide title wraps to two lines."""
        total_w = sum(int(w) for w in col_w)
        cx = x
        hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, total_w,
                                     Inches(0.4))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = ColorsV2.GREEN
        hdr.line.fill.background()
        for j, (htxt, w) in enumerate(zip(headers, col_w)):
            tb = self.textbox(slide, cx + Inches(0.12), y + Inches(0.02),
                              w - Inches(0.2), Inches(0.36))
            p = tb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            self.run(p.add_run(), htxt, ColorsV2.TBL_HDR_T, hdr_size,
                     bold=True)
            cx += w
        for i, row in enumerate(rows):
            ry = y + Inches(0.4) + i * row_h
            rbg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ry,
                                         total_w, row_h)
            rbg.fill.solid()
            rbg.fill.fore_color.rgb = (ColorsV2.ROW_A if i % 2 == 0
                                       else ColorsV2.ROW_B)
            rbg.line.color.rgb = ColorsV2.WHITE
            rbg.line.width = Pt(0.75)
            cx = x
            for j, (cell, w) in enumerate(zip(row, col_w)):
                tb = self.textbox(slide, cx + Inches(0.12),
                                  ry + Inches(0.03), w - Inches(0.2),
                                  row_h - Inches(0.05))
                tf = tb.text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
                self.run(p.add_run(), cell, ColorsV2.TBL_TXT, row_size,
                         bold=(j == 0 and first_col_bold))
                cx += w
        return y + Inches(0.4) + len(rows) * row_h

    def fig_ph(self, slide, x, y, w, h, label):
        """Light figure placeholder — matches white R-plot figures."""
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = ColorsV2.PH_BG
        box.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        box.line.width = Pt(1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self.run(p.add_run(), f"[ INSERT: {label} ]", ColorsV2.PH_TXT, 13,
                 italic=True)
        return box

    def footnote(self, slide, text, y=Inches(5.85)):
        """Small dim footnote / image-credit line at the bottom."""
        tb = self.textbox(slide, Inches(0.55), y, Inches(8.9), Inches(0.35))
        p = tb.text_frame.paragraphs[0]
        self.run(p.add_run(), text, ColorsV2.DIM, 10.5)
        return tb

    # ── Save ─────────────────────────────────────────────────

    def save(self, path):
        """Save the presentation and print summary."""
        self.pres.save(path)
        n = len(self.pres.slides)
        print(f"Saved {n} slides to: {path}")
        return path
